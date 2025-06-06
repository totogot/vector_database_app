import pandas as pd
import json
import hashlib
import numpy as np
import fitz
from pathlib import Path
from openai import OpenAI
import base64
import time
import os
import re
import torch
from transformers import CLIPModel, CLIPProcessor, BlipProcessor, BlipForConditionalGeneration, pipeline
from sentence_transformers import SentenceTransformer
from PIL import Image
from io import BytesIO
from pathlib import Path
from sklearn.metrics.pairwise import cosine_similarity
from datetime import datetime
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io



class VectorDatabase:
    def __init__(
        self,
        text_embedding_model: str,
        image_embedding_model: str,
        response_model: str,
        captioning_model: str = None, 
        openai_api_key: str = None,
        huggingface_key: str = None,
        save_dir: str = None
        ):

        # run conditional checks on model arguments passed
        valid_text_models = ["local-bge-base-en", "openai-text-embedding-3-small", "openai-text-embedding-3-large"]
        valid_image_models = ["local-clip-vit-base-patch32", "local-clip-vit-large-patch14"]
        valid_response_models = ["local-mistral-3", "openai-gpt-4o"]
        valid_caption_models = [None, "local-blip-2", "openai-gpt-4v"]

        if text_embedding_model not in valid_text_models:
            raise ValueError(f"text_embedding_model must be one of {valid_text_models}")
        if image_embedding_model not in valid_image_models:
            raise ValueError(f"image_embedding_model must be one of {valid_image_models}")
        if response_model not in valid_response_models:
            raise ValueError(f"response_model must be one of {valid_response_models}")
        if captioning_model not in valid_caption_models:
            raise ValueError(f"captioning_model must be one of {valid_caption_models}")

        # check OpenAI key has been provided for the API related models
        openai_related_models = [model for model in [text_embedding_model, image_embedding_model, captioning_model] if "openai" in model]
        if openai_related_models and not openai_api_key:
            raise ValueError("OpenAI API key is required for selected models")

        # if key provided activate client
        if openai_api_key:
            self.openai_client = OpenAI(api_key=openai_api_key)
        
        # check if huggingface key has been provided for the gated models
        if valid_response_models == "local-mistral-3" and not huggingface_key:
            raise ValueError("Huggingface token is required for accessing gated models")
        
        if huggingface_key: 
            self.huggingface_key = huggingface_key

        # set up the embedding and captioning functions

        # text embedding
        text_embeddings_route = {
            "local-bge-base-en": self.bge_text_embedder, 
            "openai-text-embedding-3-small": self.openai_text_embedder, 
            "openai-text-embedding-3-large": self.openai_text_embedder
        }

        # image embedding
        image_embeddings_route = {
            "local-clip-vit-base-patch32": self.clip_base_image_embedder, 
            "local-clip-vit-large-patch14": self.clip_large_image_embedder
        }

        # captioning
        image_captionings_route = {
            "local-blip-2": self.blip_caption_image, 
            "openai-gpt-4v": self.openai_caption_image
        }

        # response generation
        summary_response_route = {
            "local-mistral-3": self.generate_mistral_response, 
            "openai-gpt-4o": self.generate_gpt_response
        }

        self.text_embedding_model = text_embedding_model
        self.image_embedding_model = image_embedding_model
        self.captioning_model = captioning_model if captioning_model is not None else None

        self.text_embedding_function = text_embeddings_route.get(text_embedding_model)
        self.image_embedding_function = image_embeddings_route.get(image_embedding_model)
        self.summary_response_function = summary_response_route.get(response_model)
        self.image_captioning_function = image_captionings_route.get(captioning_model) if captioning_model is not None else None

        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        
        # load models to device
        if text_embedding_model == "local-bge-base-en":
            print("Loading BGE embedder")
            self.bge_embedding_model = SentenceTransformer("BAAI/bge-base-en").to(self.device)

        if image_embedding_model == "local-clip-vit-base-patch32":
            print("Loading CLIP base embedder")
            self.clip_base_image_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
            self.clip_base_image_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")

        elif image_embedding_model == "local-clip-vit-large-patch14":
            print("Loading CLIP large embedder")
            self.clip_large_image_model = CLIPModel.from_pretrained("openai/clip-vit-large-patch14").to(self.device)
            self.clip_large_image_processor = CLIPProcessor.from_pretrained("openai/clip-vit-large-patch14")

        if captioning_model == "local-blip-2":
            print("Loading BLIP captioning model")
            self.blip_captioning_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base").to(self.device)
            self.blip_captioning_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")

        if response_model == "local-mistral-3":
            print("Loading Mistral response model")
            self.mistral_response_model = pipe = pipeline(
                "image-text-to-text", 
                model="mistralai/Mistral-Small-3.1-24B-Instruct-2503", 
                device=self.device,
                token=self.huggingface_key,
                torch_dtype=torch.bfloat16
            )



        # define the directory for the vector database
        script_dir = os.path.dirname(os.path.abspath(__file__))
        self.default_db_folder = "vector_db"
        self.vector_db_folder = os.path.join(script_dir, self.default_db_folder) if save_dir is None else os.path.join(script_dir, save_dir)

        os.makedirs(self.vector_db_folder, exist_ok=True)

        self.text_data = self._load_pickle(os.path.join(self.vector_db_folder, "text_data.pkl"))
        self.image_data = self._load_pickle(os.path.join(self.vector_db_folder, "image_data.pkl"))

        self.search_granularity = None

        
    
    #############################################

    ####### FUNCTIONS FOR OPERATING DATABASE

    #############################################

    def _load_pickle(self, path):
        return pd.read_pickle(path) if os.path.exists(path) else pd.DataFrame()
    
    def vectorize_folder(self, folder_path):
        folder_files = list(Path(folder_path).rglob("*"))
        all_files = [f for f in folder_files if f.is_file()]

        for file in all_files:
            # Vectorize each file in the folder - but don't save individually
            try:
                self.vectorize_file(file, save_file_vectorizer=False)
            except Exception as e:
                print(f"Error processing file {file}: {e}")
                continue
        
        # Save after vectorizing all files
        self._save_vector_db()

        return

    def vectorize_file(self, file_path, save_file_vectorizer=True):
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1]
        file_hash = self.get_file_hash(file_path)



        # run conditional checks to ensure file is not already processed
        if self.file_already_processed(file_hash):
            print(f"File {file_name} already processed - skipping")
            return
        else:
            print(f"Processing file: {file_name}")
            self.file_hash = file_hash
            
        if ext == ".pdf":
            print("PDF detected")
            file_text, file_images = self.embed_pdf(file_path)
            
            self.text_data = pd.concat(
                [self.text_data, file_text],
                ignore_index=True
            )
            self.image_data = pd.concat(
                [self.image_data, file_images],
                ignore_index=True
            )
        
        elif ext == ".pptx":
            print("PPTX detected")
            file_text, file_images = self.embed_pptx(file_path)
            
            self.text_data = pd.concat(
                [self.text_data, file_text],
                ignore_index=True
            )
            self.image_data = pd.concat(
                [self.image_data, file_images],
                ignore_index=True
            )

        else:
            print(f"File type {ext} not supported - currently only PDF and PPTX file types supported")
            return
        
        # If vectorize called on single file then save directly
        if save_file_vectorizer:
            self._save_vector_db()

        return
    
    def _save_vector_db(self):
        if not os.path.exists(self.vector_db_folder):
            os.makedirs(self.vector_db_folder, exist_ok=True)

        self.text_data.to_pickle(os.path.join(self.vector_db_folder, "text_data.pkl"))
        self.image_data.to_pickle(os.path.join(self.vector_db_folder, "image_data.pkl"))

        return  
    

    #############################################

    ####### FUNCTIONS FOR PARSING FILES

    #############################################

    def embed_pdf(self, file):
        timestamp = self.get_file_timestamp(file)
        doc = fitz.open(file)

        file_text = pd.DataFrame()
        file_images = pd.DataFrame()

        for page_num, page in enumerate(doc):

            # Extract text and bounding boxes
            text_blocks = page.get_text("blocks")
            for idx, block in enumerate(text_blocks):
                x0, y0, x1, y1, text = block[0:5]
                if text.strip():
                    txt_entry = pd.DataFrame(
                        [{
                            "doc_name": Path(file),
                            "doc_type": "pdf",
                            "page_num": page_num,
                            "content_type": "text_chunk",
                            "content_id": f"{idx}",
                            "content_raw": text.strip(),
                            "embedding": self.text_embedding_function(text.strip()),
                            "file_hash": self.file_hash,
                            "timestamp": timestamp,
                            "bbox": [x0, y0, x1, y1]
                        }], 
                        index = [0]
                    )

                    file_text = pd.concat(
                        [file_text, txt_entry], 
                        ignore_index=True
                    )

            # Extract images and bounding boxes
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = doc.extract_image(xref)
                base64_image = base64.b64encode(base_image["image"]).decode("utf-8")
                bbox = page.get_image_bbox(img[7], transform=False) #

                img_entry = pd.DataFrame(
                    [{
                        "doc_name": Path(file),
                        "doc_type": "pdf",
                        "page_num": page_num,
                        "content_type": "image",
                        "content_id": f"{xref}",
                        "content_raw": base64_image,
                        "embedding": self.image_embedding_function(base64_image),
                        "file_hash": self.file_hash,
                        "timestamp": timestamp,
                        "bbox": [bbox.x0, bbox.y0, bbox.x1, bbox.y1]
                    }],
                    index = [0]
                )

                file_images = pd.concat(
                        [file_images, img_entry], 
                        ignore_index=True
                    )

                if self.image_captioning_function is not None:
                    print(f"Captioning Image: pg {page_num}; img {xref}")
                    max_retries = 3
                    for attempt in range(1, max_retries +1): # we will attempt to generate 3 times
                        try:
                            caption = self.image_captioning_function(base64_image)
                            txt_entry = pd.DataFrame(
                                [{
                                    "doc_name": Path(file),
                                    "doc_type": "pdf",
                                    "page_num": page_num,
                                    "content_type": "image_caption",
                                    "content_id": f"{xref}",
                                    "content_raw": caption.strip(),
                                    "embedding": self.text_embedding_function(caption.strip()),
                                    "file_hash": self.file_hash,
                                    "timestamp": timestamp,
                                    "bbox": [bbox.x0, bbox.y0, bbox.x1, bbox.y1]
                                }], 
                                index = [0]
                            )

                            file_text = pd.concat(
                                [file_text, txt_entry], 
                                ignore_index=True
                            )
                            break
                        except Exception as e:
                            print(f"Attempt {attempt} failed for Page {page_num}, Img {xref}: {e}")
                            if attempt < max_retries:
                                time.sleep(1)

        return file_text, file_images

    def embed_pptx(self, file):
        timestamp = self.get_file_timestamp(file)
        prs = Presentation(file)

        file_text = pd.DataFrame()
        file_images = pd.DataFrame()

        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:

                # process text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    txt_entry = pd.DataFrame([{
                        "doc_name": Path(file),
                        "doc_type": "pptx",
                        "page_num": slide_num,
                        "content_type": "text_chunk",
                        "content_id": f"{shape.shape_id}",
                        "content_raw": text,
                        "embedding": self.text_embedding_function(text),
                        "file_hash": self.file_hash,
                        "timestamp": timestamp,
                        "bbox": [shape.left, shape.top, shape.left + shape.width, shape.top + shape.height]
                    }])
                    file_text = pd.concat([file_text, txt_entry], ignore_index=True)

                # process image shapes
                if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        base64_image = base64.b64encode(image_bytes).decode("utf-8")

                        img_entry = pd.DataFrame([{
                            "doc_name": Path(file),
                            "doc_type": "pptx",
                            "page_num": slide_num,
                            "content_type": "image",
                            "content_id": f"{shape.shape_id}",
                            "content_raw": base64_image,
                            "embedding": self.image_embedding_function(base64_image),
                            "file_hash": self.file_hash,
                            "timestamp": timestamp,
                            "bbox": [shape.left, shape.top, shape.left + shape.width, shape.top + shape.height]
                        }])
                        file_images = pd.concat([file_images, img_entry], ignore_index=True)

                        if self.image_captioning_function is not None:
                            print(f"Captioning Image: slide {slide_num}; shape {shape.shape_id}")
                            for attempt in range(3):
                                try:
                                    caption = self.image_captioning_function(base64_image)
                                    txt_entry = pd.DataFrame([{
                                        "doc_name": Path(file),
                                        "doc_type": "pptx",
                                        "page_num": slide_num,
                                        "content_type": "image_caption",
                                        "content_id": f"{shape.shape_id}",
                                        "content_raw": caption.strip(),
                                        "embedding": self.text_embedding_function(caption.strip()),
                                        "file_hash": self.file_hash,
                                        "timestamp": timestamp,
                                        "bbox": [shape.left, shape.top, shape.left + shape.width, shape.top + shape.height]
                                    }])
                                    file_text = pd.concat([file_text, txt_entry], ignore_index=True)
                                    break
                                except Exception as e:
                                    print(f"Attempt {attempt+1} failed for Slide {slide_num}, Shape {shape.shape_id}: {e}")
                                    if attempt < 2:
                                        time.sleep(1)
                else:
                    print(f"Shape {shape.shape_id} on slide {slide_num} is a picture but has no image data.")

        return file_text, file_images

    @staticmethod
    def get_file_hash(file_path):
        """
        Returns hash of file to determine if aything has changed.
        """
        hasher = hashlib.sha256()
        with open(file_path, 'rb') as f:
            buf = f.read()
            hasher.update(buf)
        return hasher.hexdigest()
    
    def file_already_processed(self, file_hash):
        text_check = "file_hash" in self.text_data.columns and self.text_data["file_hash"].isin([file_hash]).any()
        image_check = "file_hash" in self.image_data.columns and self.image_data["file_hash"].isin([file_hash]).any()
        return text_check or image_check
    
    @staticmethod
    def get_file_timestamp(file_path):
        """
        Returns the last modified timestamp of the given file as an ISO formatted string.
        """
        timestamp = os.path.getmtime(file_path)
        return datetime.fromtimestamp(timestamp).isoformat()


    #############################################

    ####### FUNCTIONS FOR EMBEDDING DATA

    #############################################

    def openai_text_embedder(self, text):
        """
        Function for routing text embedding inference to OpenAI API endpoint
        """
        response = self.openai_client.embeddings.create(
            input = text,
            model = self.text_embedding_model.replace("openai-","")
        )
        embedding = np.array(response.data[0].embedding)
        embedding = embedding / np.linalg.norm(embedding)

        return embedding.squeeze()
    
    def bge_text_embedder(self, text):
        """
        Function for routing text embedding inference to local huggingface BGE model
        """
        model = self.bge_embedding_model
        embedding = model.encode(text, normalize_embeddings=True, convert_to_tensor=True, torch_dtype=torch.float16)
        
        return embedding.cpu().numpy().squeeze()

    def clip_base_image_embedder(self, base64_image):
        """
        Function for routing image embedding inference to local huggingface CLIP base model
        """
        model = self.clip_base_image_model 
        embedding_processor = self.clip_base_image_processor

        image_bytes = base64.b64decode(base64_image)
        image = Image.open(BytesIO(image_bytes)).convert("RGB")

        inputs = embedding_processor(images=image, return_tensors="pt").to(self.device)

        with torch.no_grad():
            embedding = model.get_image_features(**inputs)

        embedding = embedding.cpu().numpy() / np.linalg.norm(embedding.cpu().numpy(), axis=1, keepdims=True)

        return embedding.squeeze()
    
    def clip_base_text_embedder(self, text):
        """
        Function for routing text embedding inference to local huggingface CLIP base model
        """
        model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
        embedding_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
        
        inputs = embedding_processor(
            text=text, return_tensors="pt", padding=True, truncation=True
            ).to(self.device)

        with torch.no_grad():
            embedding = model.get_text_features(**inputs)
        
        embedding = embedding.cpu().numpy() / np.linalg.norm(embedding.cpu().numpy(), axis=1, keepdims=True)

        return embedding.squeeze()

    def clip_large_image_embedder(self, base64_image):
        """
        Function for routing image embedding inference to local huggingface CLIP large model
        """
        model = self.clip_large_image_model
        embedding_processor = self.clip_large_image_processor
        
        image_bytes = base64.b64decode(base64_image)
        image = Image.open(BytesIO(image_bytes)).convert("RGB")

        inputs = embedding_processor(images=image, return_tensors="pt").to(self.device)

        with torch.no_grad():
            embedding = model.get_image_features(**inputs)

        embedding = embedding.cpu().numpy() / np.linalg.norm(embedding.cpu().numpy(), axis=1, keepdims=True)

        return embedding.squeeze()
    
    def clip_large_text_embedder(self, text):
        """
        Function for routing text embedding inference to local huggingface CLIP large model
        """
        model = CLIPModel.from_pretrained("openai/clip-vit-large-patch14").to(self.device)
        embedding_processor = CLIPProcessor.from_pretrained("openai/clip-vit-large-patch14")
        
        inputs = embedding_processor(
            text=text, return_tensors="pt", padding=True, truncation=True
            ).to(self.device)

        with torch.no_grad():
            embedding = model.get_text_features(**inputs)
        
        embedding = embedding.cpu().numpy() / np.linalg.norm(embedding.cpu().numpy(), axis=1, keepdims=True)

        return embedding.squeeze()

    def blip_caption_image(self, base64_image):
        """
        Function for routing image captioning inference to local huggingface Blip model
        """
        model = self.blip_captioning_model
        embedding_processor = self.blip_captioning_processor
        
        image_bytes = base64.b64decode(base64_image)
        image = Image.open(BytesIO(image_bytes)).convert("RGB")

        text = "this image shows"
        inputs = embedding_processor(image, text, return_tensors="pt").to(self.device)
        response = model.generate(**inputs)
        
        completion = embedding_processor.decode(response[0], skip_special_tokens=True)

        return completion.replace("this image shows", "").strip()
    
    def openai_caption_image(self, base64_image):
        """
        Function for routing image captioning inference to OpenAI API endpoint
        """
        prompt = """
        can you caption this image with a description?
        provide the response in the following format:
        {{
            "description": "<image_description>"
        }}
        """

        completion = self.openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        { "type": "text", "text": prompt },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}",
                            },
                        },
                    ],
                }
            ],
        )
        caption = json.loads(self.clean_json_string(completion.choices[0].message.content))
        
        return caption["description"]

    @staticmethod
    def clean_json_string(json_string: str) -> str:
        json_string = json_string.strip()
        match = re.search(r'\{.*\}', json_string, re.DOTALL)
        if match:
            json_string = match.group(0)
        json_string = json_string.replace('```', '')
        
        return json_string.strip()
    
    
    #############################################

    ####### FUNCTIONS FOR SEARCHING DATABASE

    #############################################
    
    def run_search(
        self,
        search_content, # expecting {"text": "...", "images":[...]}
        search_location = None,
        top_n = 5
        ):

        # conditional checks to ensure correct search granularity provided
        if search_location is None:
            print("No specific search granularity provided - searching full database")
            self.search_loc = None
        elif os.path.isdir(search_location):
            print(f"Searching files in specified folder: {search_location}")
            self.search_loc = Path(search_location)
        elif os.path.isfile(search_location):
            print(f"Searching only in specified file: {search_location}")
            self.search_loc = Path(search_location)
        else:
            raise ValueError(f"Path does not exist - please provide valid location or leave blank for full database search: {search_location}")
        
        self.top_n = top_n

        # conditional checks to determine search modalities provided
        self.search_text = search_content.get("text", None)
        self.search_images = search_content.get("image", None)

        if self.search_text and self.search_images:
            self.search_mode = "text_image"
        elif self.search_text:
            self.search_mode = "text"
        elif self.search_images:
            self.search_mode = "image"
        else:
            raise ValueError("search_content must contain at least 'text' or 'image'")

        # execute search
        if len(self.text_data) > 0 or len(self.image_data) > 0:
            print("Commencing search of databse")
        else:
            raise ValueError("Database text and image data empty - please generate vector database first")
            
        if self.search_text is not None:
            text_search_results = self.run_text_search()
        else:
            text_search_results = pd.DataFrame()
        
        if self.search_images is not None:
            image_search_results = self.run_image_search()
        else:
            image_search_results = pd.DataFrame()
        
        combined_results = pd.concat([text_search_results, image_search_results], ignore_index=True)

        summary = self.summary_response_function(combined_results)
        sources = self.generate_source_list(combined_results)

        return {"response": summary, "sources": sources}

    def get_search_range(self, table):
        if self.search_loc is None:
            return_table = table
        elif os.path.isdir(self.search_loc):
            pattern = re.escape(str(self.search_loc))  # Escapes backslashes safely
            return_table = table[table['doc_name'].astype(str).str.contains(pattern, case=False)]
        elif os.path.isfile(self.search_loc):
            return_table = table[table['doc_name'] == self.search_loc]
        
        return return_table
    
    def return_similar(self, question_vector, search_vectors):

        question_vector_reshaped = question_vector.reshape(1, -1)

        reference_vectors = list(search_vectors['embedding'])
        reference_vectors = np.array(reference_vectors).reshape(len(reference_vectors), -1)

        similarities = cosine_similarity(question_vector_reshaped, reference_vectors).flatten()

        top_indices = similarities.argsort()[-self.top_n:][::-1]
        print(f"Top Similarity Scores: {np.sort(similarities)[-self.top_n:][::-1]}")

        return search_vectors.iloc[top_indices]

    def run_text_search(self):
        text = self.search_text
        text_references = pd.DataFrame()
        image_references = pd.DataFrame()

        # search text vs text
        if len(self.text_data) > 0:
            embedded_search_text = self.text_embedding_function(text.strip())
            search_range = self.get_search_range(self.text_data)
            try:
                print(f"Searching text vs text - returned range len: {len(search_range)}")
                text_references = self.return_similar(embedded_search_text, search_range)
                text_references['search_reference'] = text
            except Exception as e:
                print(f"Searching text vs text - returned range len: {len(search_range)}")
                text_references = pd.DataFrame()
        
        # search text vs image
        text_image_embedding = {
            "local-clip-vit-base-patch32": self.clip_base_text_embedder, 
            "local-clip-vit-large-patch14": self.clip_large_text_embedder
        }

        if len(self.image_data) > 0:
            embedded_search_text = text_image_embedding.get(self.image_embedding_model)(text)
            search_range = self.get_search_range(self.image_data)
            try:
                print(f"Searching text vs images - returned range len: {len(search_range)}")
                image_references = self.return_similar(embedded_search_text, search_range)
                image_references['search_reference'] = text
            except Exception as e:
                print(f"Searching text vs images - returned range len: {len(search_range)}")
                image_references = pd.DataFrame()
        
        # combine results 
        combined_references = pd.concat([text_references, image_references], ignore_index=True)
        combined_references[['doc_name', 'page_num','content_type', 'content_id', 'content_raw']]

        return combined_references
            
    def run_image_search(self):
        images = self.search_images
        image_references = pd.DataFrame()

        # search image vs image
        if len(self.image_data) > 0:
            search_range = self.get_search_range(self.image_data)
            try:
                print(f"Searching images vs images - returned range len: {len(search_range)}")
                for img in images:
                    embedded_search_image = self.image_embedding_function(img)
                    single_image_references = self.return_similar(embedded_search_image, search_range)
                    single_image_references['search_reference'] = img
                    image_references = pd.concat([image_references, single_image_references], ignore_index=True)
            except Exception as e:
                print(f"Searching images vs images - returned range len: {len(search_range)}")
                image_references = pd.DataFrame()
        
        # search image vs text
        ##### TO DO: Add image vs text search #####

        return image_references
    

    #############################################

    ####### FUNCTIONS FOR GENERATING RESPONSES

    #############################################

    def generate_mistral_response(self, response):

        # initiate the pipeline
        pipe = self.mistral_response_model
        # pipeline(
        #     "image-text-to-text", 
        #     model="mistralai/Mistral-Small-3.1-24B-Instruct-2503", 
        #     device=self.device,
        #     token=self.huggingface_key,
        #     torch_dtype=torch.bfloat16
        #     )

        # define the query material
        content_list = [
            {"type": "text", "text": "Here is the search query content:\n"}
            ]
        
        if self.search_text:
            content_list.append({"type": "text", "text": f"Text query: {self.search_text}"})
        if self.search_images:
            content_list.append({"type": "image", "url": f"data:image/jpeg;base64,{self.search_images}"})

        # combine relevant retrieval content - from the response DataFrame
        content_list.append({"type": "text", "text": "\n\nHere is the retreived material:"})

        # append any text content
        text_content = response[response["content_type"] != "image"]
        for row in text_content.itertuples():
            content_list.append({
                "type": "text",
                "text": row.content_raw,
            })

        # append any image content
        image_content = response[response["content_type"] == "image"]
        for row in image_content.itertuples():
            content_list.append({
                "type": "image",
                "url": f"data:image/jpeg;base64,{row.content_raw}",
            })

        # append a closing instruction
        content_list.append(
            {
                "type": "text", 
                "text": "Please generate an appropriate response to the query based on the information available - you can choose what retrieved information is most appropriate to answer the question."
            }
        )

        # run pipeline inference
        try:
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful assistant that is going to summarize the key points from search results, based on a user query and relevant retrieved content."
                },
                {
                    "role": "user",
                    "content": content_list
                }
            ]

            outputs = pipe(text=messages, max_new_tokens=100, return_full_text=False)

            # Extract and return the Mistral response
            return outputs[0]["generated_text"]
        except Exception as e:
            print(f"Error running Mistral local inference: {e}")
            return None
        
    def generate_gpt_response(self, response):

        # define the query material
        content_list = [
            {"type": "text", "text": "Here is the search query content:\n"}
            ]
        
        if self.search_text:
            content_list.append({"type": "text", "text": f"Text query: {self.search_text}"})
        if self.search_images:
            print(self.search_images)
            for image in self.search_images:
                content_list.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image}"}})

        # combine relevant retrieval content - from the response DataFrame
        content_list.append({"type": "text", "text": "\n\nHere is the retreived material:"})

        # append any text content
        text_content = response[response["content_type"] != "image"]
        for row in text_content.itertuples():
            content_list.append({
                "type": "text",
                "text": row.content_raw,
            })

        # append any image content
        image_content = response[response["content_type"] == "image"]
        for row in image_content.itertuples():
            content_list.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{row.content_raw}",
                },
            })

        # append a closing instruction
        content_list.append(
            {
                "type": "text", 
                "text": "Please generate an appropriate response to the query based on the information available - you can choose what retrieved information is most appropriate to answer the question."
            }
        )

        # Call the GPT endpoint
        try:
            completion = self.openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": "You are a helpful assistant that is going to summarize the key points from search results, based on a user query and relevant retrieved content."
                    },
                    {
                        "role": "user",
                        "content": content_list
                    }
                ]
            )

            # Extract and return the GPT response
            return completion.choices[0].message.content.strip()
        except Exception as e:
            print(f"Error calling GPT endpoint: {e}")
            return None
        
    def generate_source_list(self, response):

        sources = []
        for row in response.itertuples():
            sources.append(
                {
                    "doc_name": str(row.doc_name),
                    "page_num": row.page_num,
                    "content_type": row.content_type,
                    "content_id": row.content_id,
                    "content_raw": row.content_raw
                }
            )
        return sources
